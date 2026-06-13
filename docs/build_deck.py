# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY=RGBColor(0x0d,0x21,0x37); BLUE=RGBColor(0x1f,0x4e,0x79); ACC=RGBColor(0x2f,0x6f,0xb5)
GREEN=RGBColor(0x1b,0x5e,0x20); AMBER=RGBColor(0xb8,0x86,0x00); RED=RGBColor(0xb7,0x1c,0x1c)
GREY=RGBColor(0x55,0x60,0x70); LIGHT=RGBColor(0xf0,0xf5,0xfb); WHITE=RGBColor(0xff,0xff,0xff)
FONT="Tahoma"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]

def box(s,l,t,w,h,fill=None,line=None):
    sp=s.shapes.add_shape(1,l,t,w,h); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    return sp

def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=6,line_sp=1.06):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0); p.line_spacing=line_sp
        if isinstance(para,tuple): para=[para]
        for seg in para:
            text,sz,col,bold=(seg+(False,))[:4] if len(seg)==3 else seg
            r=p.add_run(); r.text=text; f=r.font; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=FONT
    return tb

def bg(s,color=WHITE):
    box(s,0,0,SW,SH,fill=color)

def header(s,title,kicker=None,color=NAVY):
    box(s,0,0,SW,Inches(1.15),fill=color)
    box(s,0,Inches(1.15),SW,Pt(4),fill=ACC)
    if kicker: txt(s,Inches(0.55),Inches(0.16),Inches(11),Inches(0.3),[[(kicker,12,RGBColor(0x9f,0xc0,0xe8),True)]])
    txt(s,Inches(0.55),Inches(0.40),Inches(12.2),Inches(0.7),[[(title,26,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)

def footer(s,n):
    txt(s,Inches(0.55),Inches(7.05),Inches(8),Inches(0.3),[[("PWMS · PSA-HKT Workforce Management System",10,GREY)]])
    txt(s,Inches(12.0),Inches(7.05),Inches(0.9),Inches(0.3),[[(str(n),10,GREY,True)]],align=PP_ALIGN.RIGHT)

def content(title,kicker,blocks,n,foot=True):
    s=prs.slides.add_slide(blank); bg(s); header(s,title,kicker)
    y=Inches(1.5)
    for b in blocks:
        runs=b["runs"]; h=b.get("h",Inches(0.9)); card=b.get("card",False); col=b.get("col",LIGHT)
        if card:
            box(s,Inches(0.55),y,Inches(12.23),h,fill=col)
            bar=b.get("bar"); 
            if bar: box(s,Inches(0.55),y,Pt(6),h,fill=bar)
            txt(s,Inches(0.8),y+Pt(6),Inches(11.8),h-Pt(12),runs,anchor=MSO_ANCHOR.MIDDLE)
        else:
            txt(s,Inches(0.7),y,Inches(12.0),h,runs)
        y=y+h+Inches(0.12)
    if foot: footer(s,n)
    return s

# ---------- 1 COVER ----------
s=prs.slides.add_slide(blank); bg(s,NAVY)
box(s,0,Inches(4.7),SW,Pt(5),fill=ACC)
txt(s,Inches(0.8),Inches(1.0),Inches(11.7),Inches(0.4),[[("PASSENGER SERVICES · AOTGA — HKT",14,RGBColor(0x9f,0xc0,0xe8),True)]])
txt(s,Inches(0.8),Inches(2.0),Inches(12),Inches(2.2),[
  [("PWMS",48,WHITE,True)],
  [("PSA-HKT Workforce Management System",24,RGBColor(0xcf,0xe0,0xf2),False)],
])
txt(s,Inches(0.8),Inches(5.0),Inches(12),Inches(1.6),[
  [("แดชบอร์ดเว็บ อ่านชีต Roster/Flight อัตโนมัติ → เช็คคนครบตาม SLA → หาคนช่วย → สร้างข้อความ SOS",16,RGBColor(0xcf,0xe0,0xf2))],
  [("เอกสารนำเสนอ · คู่มือใช้งาน · แผนพัฒนาต่อ   |   มิ.ย. 2026",13,RGBColor(0x9f,0xc0,0xe8))],
])

# ---------- 2 AGENDA ----------
content("สิ่งที่จะนำเสนอ","ภาพรวม",[
 {"card":True,"bar":BLUE,"h":Inches(1.1),"runs":[[("1 · เสนอระบบ",16,BLUE,True),("   ที่มา ปัญหาเดิม ภาพรวมระบบ และ data flow",13,GREY)]]},
 {"card":True,"bar":GREEN,"h":Inches(1.1),"runs":[[("2 · คู่มือ / วิธีใช้งาน",16,GREEN,True),("   แต่ละแท็บทำอะไร ใช้ยังไง ทีละขั้น",13,GREY)]]},
 {"card":True,"bar":AMBER,"h":Inches(1.1),"runs":[[("3 · แผนพัฒนาต่อ",16,AMBER,True),("   สิ่งที่ทำเสร็จแล้ว · กำลังจะทำ · อนาคต",13,GREY)]]},
],2)

# ======== SECTION 1 ========
def divider(num,title,sub):
    s=prs.slides.add_slide(blank); bg(s,NAVY)
    box(s,Inches(0.8),Inches(2.6),Inches(1.6),Inches(1.6),fill=ACC)
    txt(s,Inches(0.8),Inches(2.6),Inches(1.6),Inches(1.6),[[(num,54,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(2.8),Inches(2.7),Inches(9.8),Inches(1.6),[
      [(title,34,WHITE,True)],[(sub,16,RGBColor(0x9f,0xc0,0xe8))]],anchor=MSO_ANCHOR.MIDDLE)
divider("1","เสนอระบบ","System Overview")

content("ปัญหาเดิม — ทำไมต้องมีระบบนี้","ที่มา",[
 {"card":True,"bar":RED,"h":Inches(1.0),"runs":[[("ขอคนซัพแบบ manual",15,RED,True),("  แต่ละทีมพิมพ์ขอคนช่วยทางไลน์ ไม่มีมาตรฐาน หาย/ตกหล่น",13,GREY)]]},
 {"card":True,"bar":RED,"h":Inches(1.0),"runs":[[("หาคนว่างยาก",15,RED,True),("  ต้องไล่ดูตารางทีละทีมว่าใครว่าง รู้ระบบเช็คอินตรงไหม เวลาชนกันไหม",13,GREY)]]},
 {"card":True,"bar":RED,"h":Inches(1.0),"runs":[[("เช็ค SLA ด้วยมือ",15,RED,True),("  ต้องนับเองว่าแต่ละไฟลท์ส่งคนครบตามมาตรฐานสายการบินไหม",13,GREY)]]},
 {"card":True,"bar":RED,"h":Inches(1.0),"runs":[[("OT / เวลากะ ไม่ถูกตรวจ",15,RED,True),("  คนถูก assign ไฟลท์นอกเวลากะ หรือ OT เกินจำเป็น มองไม่เห็น",13,GREY)]]},
],4)

content("ระบบนี้คืออะไร","แนวคิด",[
 {"runs":[[("เว็บแดชบอร์ดตัวเดียว ",18,NAVY,True),("ที่อ่านข้อมูลจาก Google Sheets ของจริง แล้วช่วยงานหัวหน้า/ดิวตี้:",15,GREY)]],"h":Inches(0.7)},
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("อ่านอัตโนมัติ",15,BLUE,True),("  ดึง Roster (เวร) + Flight (ไฟลท์/เวลา) + Master (อัตรากำลัง) ทุกครั้งที่เปิด",13,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("เช็ค SLA",15,BLUE,True),("  เทียบคนที่ส่ง vs มาตรฐานสายการบิน (SUP/Check-in/Gate/Arrival) ต่อไฟลท์",13,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("หาคนช่วย + SOS",15,BLUE,True),("  แนะนำคนทีมอื่นที่ว่าง+รู้ระบบ มาเสริม และสร้างข้อความ SOS คัดลอกส่งไลน์",13,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("จัดเวรอัตโนมัติ",15,BLUE,True),("  เติมคนข้ามทีม หรือจัดเวรใหม่ทั้งวันให้ครบตาม SLA",13,GREY)]]},
],5)

content("ภาพรวม 8 แท็บการทำงาน","เมนู",[
 {"card":True,"col":LIGHT,"h":Inches(0.72),"runs":[[("▦ Dashboard",14,NAVY,True),("   สรุปยอดคนเข้าเวร/หยุด/ลา · OT รวม",12,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(0.72),"runs":[[("✈️ Flights & SLA",14,NAVY,True),("   ทุกไฟลท์ + คนที่ต้องการ vs ส่งจริง + สถานะครบ/ขาด",12,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(0.72),"runs":[[("🆘 Support",14,RED,True),("   ไฟลท์คนไม่ครบ + แนะนำคนช่วย + ข้อความ SOS",12,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(0.72),"runs":[[("🧭 ตรวจ Assign",14,NAVY,True),("   เช็คคนรายคน: เวลากะครอบไฟลท์ไหม · OT · ช่วงว่าง",12,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(0.72),"runs":[[("🤖 เติม / Auto Assign",14,NAVY,True),("   เติมคนข้ามทีม · จัดเวรใหม่ทั้งวันตาม SLA",12,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(0.72),"runs":[[("📅 จัดล่วงหน้า · ⏱️ OT · ☰ Timetable",14,NAVY,True),("   วางแผนล่วงหน้า+export ต่อทีม · OT รายเดือน · ตารางรายคน",12,GREY)]]},
],6)

content("Data Flow — ข้อมูลไหลยังไง","สถาปัตยกรรม",[
 {"card":True,"bar":ACC,"h":Inches(1.05),"runs":[[("1) แหล่งข้อมูล (Google Sheets)",15,BLUE,True)],[("ROSTER (เวร/ตำแหน่ง/ไฟลท์ที่ลง) · FLIGHT (STA/STD) · MASTER (อัตรากำลัง+สถานะ Active/Resigned)",12,GREY)]]},
 {"card":True,"bar":ACC,"h":Inches(1.05),"runs":[[("2) ตัวอ่าน + เครื่องคำนวณ (Apps Script)",15,BLUE,True)],[("อ่านแถวแบบ dynamic → จับคู่คน-ไฟลท์ → เทียบ SLA รายสาย → หาคนว่าง/รู้ระบบ",12,GREY)]]},
 {"card":True,"bar":ACC,"h":Inches(1.05),"runs":[[("3) แสดงผล (เว็บแดชบอร์ด)",15,BLUE,True)],[("8 แท็บ · เลือกวันที่ · ข้อความ SOS · export ชีตต่อทีม",12,GREY)]]},
 {"card":True,"col":RGBColor(0xe8,0xf5,0xe9),"h":Inches(0.8),"runs":[[("✓ อ่านชีตสดทุกครั้ง → เพิ่ม/ลาออกคน หรือแก้ไฟลท์ ระบบเห็นทันที ไม่ต้องแก้โค้ด",13,GREEN,True)]]},
],7)

# ======== SECTION 2 ========
divider("2","คู่มือ / วิธีใช้งาน","User Manual")

content("เริ่มใช้งาน","ขั้นพื้นฐาน",[
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("1) เปิดเว็บแอป",15,BLUE,True),("   เปิดลิงก์เว็บแอป (จาก Apps Script) บนคอม/มือถือ",13,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("2) เลือกวันที่",15,BLUE,True),("   ระบบจะดึง roster ของวันนั้นมาคำนวณให้",13,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.95),"runs":[[("3) เลือกแท็บที่ต้องการ",15,BLUE,True),("   แต่ละแท็บโหลดเมื่อกด (เร็วขึ้น) · มีช่องค้นหา + เลือกทีมในทุกตาราง",13,GREY)]]},
 {"card":True,"col":RGBColor(0xff,0xf8,0xe1),"h":Inches(0.85),"runs":[[("💡 ข้อมูลอัปเดตจากชีตทุกครั้งที่รีเฟรช — แก้ในชีต Google แล้วรีเฟรชเว็บ เห็นผลทันที",13,AMBER,True)]]},
],9)

content("✈️ Flights & SLA — อ่านยังไง","คู่มือแท็บ",[
 {"runs":[[("แต่ละแถว = 1 ไฟลท์ · คอลัมน์ ",14,NAVY,True),("จัด/รวม",13,BLUE,True),(" = คนที่ส่งจริง / ที่ต้องการ · ช่อง SUP·Check-in·Gate·Arrival = แยกตามหน้าที่",13,GREY)]],"h":Inches(0.8)},
 {"card":True,"col":RGBColor(0xe8,0xf5,0xe9),"h":Inches(0.7),"runs":[[("✅ ครบ",14,GREEN,True),("   ส่งคนครบตามมาตรฐานสายการบินนั้น",13,GREY)]]},
 {"card":True,"col":RGBColor(0xfd,0xec,0xec),"h":Inches(0.7),"runs":[[("⚠️ ขาด …",14,RED,True),("   บอกว่าขาดหน้าที่ไหน กี่คน (เช่น Gate ขาด 2)",13,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(1.3),"runs":[
   [("กฎสำคัญที่ระบบใช้:",13,NAVY,True)],
   [("• leg-based: ไฟลท์ขาออกอย่างเดียวไม่นับ Arrival · ขาเข้าอย่างเดียวไม่นับ Check-in/Gate",12,GREY)],
   [("• 00:00 = ค่า placeholder (ยังไม่กรอกเวลา) ไม่ใช่เที่ยงคืนจริง",12,GREY)],
   [("• ไฟลท์ไม่มี STA/STD เลย → ขึ้นเตือน “เติมเวลาในชีต”",12,GREY)]]},
],10)

content("🆘 Support — หาคนช่วย + SOS","คู่มือแท็บ",[
 {"card":True,"bar":RED,"h":Inches(0.8),"runs":[[("รายการไฟลท์ที่คนไม่ครบ",14,RED,True),("   แยกตามหน้าที่ที่ขาด + ช่วงเวลา + ระบบเช็คอินที่ต้องใช้",12,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.8),"runs":[[("แนะนำคนช่วย (ทีมอื่นเท่านั้น)",14,BLUE,True),("   ว่างช่วงนั้น + รู้ระบบ · เลือกทีมที่จะดึง · 1 ช่อง = 1 คน",12,GREY)]]},
 {"card":True,"bar":AMBER,"h":Inches(0.8),"runs":[[("OFF · re-sked",14,AMBER,True),("   ถ้าคนเข้าเวรไม่พอ ระบบเสนอคน “วันหยุด” มา re-sked (ต่อท้ายรายการ)",12,GREY)]]},
 {"card":True,"col":RGBColor(0xe8,0xf5,0xe9),"h":Inches(1.0),"runs":[[("📋 ปุ่ม “ข้อความ SOS”",14,GREEN,True),("   กดคัดลอก → วางในไลน์ได้เลย",12,GREY)],[("รูปแบบ: ชื่อไฟลท์ + เวลา + หน้าที่ที่ขาด + รายชื่อคนช่วย (ชื่อ/ทีม) — เฉพาะคนทีมอื่น",12,GREY)]]},
],11)

content("🧭 ตรวจ Assign รายคน","คู่มือแท็บ",[
 {"runs":[[("ตรวจความเหมาะสมการจัดคนทีละคน — เน้นเฉพาะที่ต้องดู (🔴/🟡)",14,NAVY,True)]],"h":Inches(0.55)},
 {"card":True,"col":RGBColor(0xfd,0xec,0xec),"h":Inches(0.75),"runs":[[("🔴 ไฟลท์นอกเวลา / OT ไม่พอ",14,RED,True),("   ไฟลท์ที่ได้รับเริ่ม/เลิกนอกเวลากะ แล้ว OT ไม่ครอบ → ควร Re-Sked/ให้ OT",12,GREY)]]},
 {"card":True,"col":RGBColor(0xff,0xf8,0xe1),"h":Inches(0.75),"runs":[[("🟡 OT อาจเกิน / มีช่วงว่างยาว",14,AMBER,True),("   ไฟลท์อยู่ในเวลากะแต่ให้ OT · หรือมีช่องว่างกลางวันนาน",12,GREY)]]},
 {"card":True,"col":RGBColor(0xe8,0xf5,0xe9),"h":Inches(0.9),"runs":[[("✓ เทรน/อบรม/ประชุม นับเป็น “งาน”",14,GREEN,True)],[("ไม่ขึ้น 🔴 ว่านอกเวลา และถือว่าคนนั้น busy (ไม่ถูกดึงไปช่วยตอนอบรม)",12,GREY)]]},
],12)

content("🤖 เติม / Auto Assign · 📅 จัดล่วงหน้า","คู่มือแท็บ",[
 {"card":True,"bar":BLUE,"h":Inches(0.85),"runs":[[("🤖 เติมจาก Assign เดิม",14,BLUE,True),("   เก็บเวรเดิมไว้ แล้วเติมคนข้ามทีมเฉพาะไฟลท์ที่ขาด",12,GREY)]]},
 {"card":True,"bar":BLUE,"h":Inches(0.85),"runs":[[("🤖 Auto Assign (จัดใหม่ทั้งหมด)",14,BLUE,True),("   ล้างแล้วจัดทุกคนลงไฟลท์ให้ครบ SLA + โชว์คนพัก/สำรอง",12,GREY)]]},
 {"card":True,"bar":GREEN,"h":Inches(0.85),"runs":[[("📅 จัดล่วงหน้า + 📤 export",14,GREEN,True),("   วางแผนวันถัดไป · สร้างไฟล์ชีตแจกทีม (1 แท็บ/ทีม) ไม่ทับไฟล์จริง",12,GREY)]]},
 {"card":True,"col":LIGHT,"h":Inches(0.85),"runs":[[("🛄 SU/SQ เช็คอินคอมมอน",14,NAVY,True),("   เคาน์เตอร์รวมหมุนเวียน ≤3 ชม./คน + เกทต่อไฟลท์ (คนต่อเนื่องจากเช็คอิน)",12,GREY)]]},
],13)

content("การดูแลข้อมูล (ไม่ต้องแก้โค้ด)","Operations",[
 {"card":True,"bar":GREEN,"h":Inches(0.95),"runs":[[("เพิ่มพนักงานใหม่",15,GREEN,True),("   เพิ่มแถวในชีต (รหัส 6-8 หลัก + ชื่อ + กะ + ตำแหน่ง) → ระบบอ่านเอง",13,GREY)]]},
 {"card":True,"bar":GREEN,"h":Inches(0.95),"runs":[[("พนักงานลาออก",15,GREEN,True),("   ใส่สถานะ Resigned ใน MASTER → ตัดออกจากอัตรากำลังอัตโนมัติ",13,GREY)]]},
 {"card":True,"bar":ACC,"h":Inches(0.95),"runs":[[("Auto-deploy",15,BLUE,True),("   แก้โค้ด → push → GitHub Action รัน clasp push ขึ้น Apps Script อัตโนมัติ",13,GREY)]]},
 {"card":True,"col":RGBColor(0xff,0xf8,0xe1),"h":Inches(0.8),"runs":[[("⚠️ ต้องรักษาฟอร์แมตชีตเดิม (ตำแหน่งคอลัมน์ รหัส ชื่อ) · ใช้คำสถานะตรงตัว Active/Resigned",12,AMBER,True)]]},
],14)

# ======== SECTION 3 ========
divider("3","แผนพัฒนาต่อ","Development Roadmap")

content("สิ่งที่ทำเสร็จแล้ว","Phase 1 — Done",[
 {"card":True,"col":RGBColor(0xe8,0xf5,0xe9),"h":Inches(4.7),"runs":[
   [("✓ SLA leg-based",14,GREEN,True),("  คิดความต้องการตามขาที่ไฟลท์มีจริง (ขาออก/ขาเข้า)",12,GREY)],
   [("✓ รวมไฟลท์ซ้ำ + ซ่อนเศษขา",14,GREEN,True),("  TK172/173+TK172, EY416 codeshare → แถวเดียว",12,GREY)],
   [("✓ ชื่อทีม/ระบบ/สายการบิน ตรงชีตจริง",14,GREEN,True),("  KC→KE, รวม WY/WK, เติม QP/IX/QZ",12,GREY)],
   [("✓ ข้อความ SOS คัดลอกส่งไลน์",14,GREEN,True),("  จัดกลุ่มต่อไฟลท์ เฉพาะคนทีมอื่น",12,GREY)],
   [("✓ ดึงคน OFF มา re-sked",14,GREEN,True),("  เสนอคนวันหยุดเมื่อคนเข้าเวรไม่พอ",12,GREY)],
   [("✓ เทรน/อบรม นับเป็นงาน + busy",14,GREEN,True),("  ไม่ flag นอกเวลา · อ่านช่วงเวลาอบรมจากข้อความ",12,GREY)],
   [("✓ แก้นับ AdminDoc · กัน 00:00 placeholder · auto-deploy",14,GREEN,True)]]},
],16)

content("กำลังจะทำต่อ (ระยะใกล้)","Phase 2 — Next",[
 {"card":True,"bar":AMBER,"h":Inches(0.92),"runs":[[("Mapping ทีม → ระบบเช็คอินที่ถนัด",14,AMBER,True),("   ให้คน OFF ของทีมที่หยุดทั้งทีม ยังขึ้นช่วยงาน Check-in ได้",12,GREY)]]},
 {"card":True,"bar":AMBER,"h":Inches(0.92),"runs":[[("อ่านเวลาอบรมทุกฟอร์แมต",14,AMBER,True),("   รองรับ 08-17, 0900-1200, 8.00-17.00 ฯลฯ ให้ busy แม่นขึ้น",12,GREY)]]},
 {"card":True,"bar":AMBER,"h":Inches(0.92),"runs":[[("Auto deployment version (/exec)",14,AMBER,True),("   ให้ลิงก์เว็บ /exec อัปเดตเองหลัง push (ไม่ต้องกด New version)",12,GREY)]]},
 {"card":True,"bar":AMBER,"h":Inches(0.92),"runs":[[("จับไฟลท์กรอกสลับหัว-ท้าย",14,AMBER,True),("   เช่น VJ808/809 vs VJ809/808 ให้รวมเป็นไฟลท์เดียว",12,GREY)]]},
],17)

content("แผนระยะยาว (อนาคต)","Phase 3 — Future",[
 {"card":True,"bar":ACC,"h":Inches(0.92),"runs":[[("บันทึก SOS ลงชีต / ส่งไลน์อัตโนมัติ",14,BLUE,True),("   เก็บประวัติคำขอ + ส่งเข้ากลุ่มไลน์ทีมโดยตรง",12,GREY)]]},
 {"card":True,"bar":ACC,"h":Inches(0.92),"runs":[[("มุมมองมือถือ (Mobile-friendly)",14,BLUE,True),("   ปรับหน้าจอให้ดิวตี้ใช้บนมือถือสะดวกขึ้น",12,GREY)]]},
 {"card":True,"bar":ACC,"h":Inches(0.92),"runs":[[("รายงาน/สถิติย้อนหลัง",14,BLUE,True),("   แนวโน้มคนขาด · OT · ทีมที่ช่วยบ่อย เพื่อวางกำลังคน",12,GREY)]]},
 {"card":True,"bar":ACC,"h":Inches(0.92),"runs":[[("คาดการณ์ความต้องการคน (Predictive)",14,BLUE,True),("   ประเมินล่วงหน้าจากตารางบิน + อัตรากำลัง",12,GREY)]]},
],18)

# ---------- CLOSE ----------
s=prs.slides.add_slide(blank); bg(s,NAVY)
box(s,0,Inches(3.0),SW,Pt(5),fill=ACC)
txt(s,Inches(0.8),Inches(2.0),Inches(12),Inches(1.0),[[("สรุป",36,WHITE,True)]])
txt(s,Inches(0.8),Inches(3.3),Inches(12),Inches(2.6),[
 [("ระบบเดียวที่ช่วยตั้งแต่ เช็ค SLA → หาคนช่วย → สร้าง SOS → จัดเวร",18,WHITE,True)],
 [("อ่านชีตสดทุกครั้ง · เพิ่ม/ลาออกคนไม่ต้องแก้โค้ด · auto-deploy",14,RGBColor(0xcf,0xe0,0xf2))],
 [("พัฒนาต่อเนื่อง: re-sked ข้ามทีม · อบรมแม่นขึ้น · รายงานย้อนหลัง",14,RGBColor(0xcf,0xe0,0xf2))],
])
txt(s,Inches(0.8),Inches(6.2),Inches(12),Inches(0.5),[[("ขอบคุณครับ · PWMS — PSA-HKT / AOTGA",14,RGBColor(0x9f,0xc0,0xe8),True)]])

out="/home/user/Bot-Assign/docs/PWMS-PSA-HKT-Workforce-Management-System.pptx"
import os; os.makedirs("/home/user/Bot-Assign/docs",exist_ok=True)
prs.save(out)
print("saved", out, "· slides:", len(prs.slides._sldIdLst))
